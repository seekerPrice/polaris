#!/usr/bin/env python3
"""Polaris pitch deck builder — VC light theme, v3 (larger type + clearer hierarchy).

Generates docs/PITCH_DECK.pptx — 15 core + 3 appendix = 18 slides — styled
with a clean VC-pitch palette (cream + indigo + semantic emerald/red/amber).

Typography is centralised in the FS_* constants block — change one value
to globally re-size that role across every slide. Section tags carry an
indigo underline rule so headers visually separate from body content.

Content + citations from docs/PITCH_DECK.md. Sources verified 2026-05-18
against public first-party data.

Usage:
    uv run python scripts/build_pitch_deck.py
"""

from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
DOCS_IMG = ROOT / "docs" / "img"
OUT = ROOT / "docs" / "PITCH_DECK.pptx"

# --- VC light palette ---------------------------------------------------
BG = RGBColor(0xFA, 0xFA, 0xFC)          # cream page bg
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUBTLE = RGBColor(0xF1, 0xF5, 0xF9)      # slate-100
BORDER = RGBColor(0xE2, 0xE8, 0xF0)      # slate-200
BORDER_STRONG = RGBColor(0xCB, 0xD5, 0xE1)  # slate-300

TEXT_PRIMARY = RGBColor(0x0F, 0x17, 0x2A)
TEXT_SECONDARY = RGBColor(0x47, 0x55, 0x69)
TEXT_MUTED = RGBColor(0x94, 0xA3, 0xB8)
TEXT_INVERSE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_INVERSE_SOFT = RGBColor(0xC7, 0xD2, 0xFE)

INDIGO = RGBColor(0x4F, 0x46, 0xE5)
INDIGO_LIGHT = RGBColor(0x63, 0x66, 0xF1)
INDIGO_DARK = RGBColor(0x37, 0x30, 0xA3)
INDIGO_SOFT = RGBColor(0xEE, 0xF2, 0xFF)

SUCCESS = RGBColor(0x05, 0x96, 0x69)
DANGER = RGBColor(0xDC, 0x26, 0x26)
WARNING = RGBColor(0xD9, 0x77, 0x06)
INFO = RGBColor(0x0E, 0xA5, 0xE9)
VIOLET = RGBColor(0x7C, 0x3A, 0xED)

FONT_SANS = "Inter"
FONT_MONO = "JetBrains Mono"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# --- Typography scale (single source of truth) --------------------------
# Display
FS_HERO = 110          # cover title "Polaris"
FS_DISPLAY = 130       # thank-you "Thank you."

# Headlines (slide titles + key call-outs)
FS_H1_LARGE = 52       # mission slide headline
FS_H1 = 44             # cover hero metric, product close-up
FS_H2_LARGE = 36       # closed-loop, unit-econ
FS_H2 = 32             # market, ask
FS_H3 = 28             # solution, competitive, business model
FS_H4 = 24             # team, A1, A3

# Subtitles / leads
FS_LEAD_LARGE = 26     # cover subline
FS_LEAD = 22           # cover subline alt, A1 lead
FS_LEAD_SM = 18        # closed loop subtitle

# Body
FS_BODY_LG = 16        # main body on product / ask slides
FS_BODY = 14           # default body
FS_BODY_SM = 13        # card body, table cells
FS_BODY_XS = 12        # tight bullets in 3-column cards

# Labels / tags / captions
FS_TAG = 13            # section eyebrow tag (was 10)
FS_TAG_SM = 11         # nested tags inside cards
FS_LABEL_BOLD = 13     # bold labels (e.g. ENGINEERING 50%)
FS_LABEL = 11          # muted captions
FS_CORNER = 11         # top-corner branding labels

# Numbers / metrics
FS_NUM_HERO = 64       # market section big numbers
FS_NUM_XL = 44         # unit-econ tile big numbers
FS_NUM_LG = 32         # market big numbers like $7.44B
FS_NUM_MD = 26         # bullet emphasis numbers
FS_NUM_SM = 18         # smaller stat numbers
FS_ICON = 22           # ✓ ✗ marks
FS_ICON_LG = 36        # step numerals in closed-loop slide

# Footer
FS_FOOTER = 10         # citation strip (was 7.5)


# --- helpers ------------------------------------------------------------

def set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(
    slide, text, left, top, width, height, *,
    font=FONT_SANS, size=FS_BODY, color=TEXT_PRIMARY, bold=False,
    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.2,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    lines = text.split("\n") if isinstance(text, str) else [str(text)]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
    return tb


def add_rect(slide, left, top, width, height, fill, *,
             line=None, line_width=Pt(0.5)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = line_width
    return shape


def add_oval(slide, left, top, width, height, fill, *,
             line=None, line_width=Pt(0.5)):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = line_width
    return shape


def add_arrow(slide, begin, end, *,
              color=INDIGO, width=Pt(1.5),
              connector_type=MSO_CONNECTOR.STRAIGHT,
              dashed=False):
    """Add a connector line with a triangle arrowhead at the END point.

    `begin` and `end` are (x, y) tuples of Emu/Inches values.
    Use MSO_CONNECTOR.ELBOW for L-shaped right-angle routing.
    """
    bx, by = begin
    ex, ey = end
    conn = slide.shapes.add_connector(connector_type, bx, by, ex, ey)
    conn.line.color.rgb = color
    conn.line.width = width
    if dashed:
        conn.line.dash_style = 7  # MSO_LINE.ROUND_DOT-ish (XML enum 7)
    # Inject <a:tailEnd> for arrowhead
    ln = conn.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")
    return conn


def add_tag(
    slide, text, left, top, *,
    color=INDIGO, size=FS_TAG, with_rule=True, rule_width=0.5,
):
    """Section eyebrow tag with optional underline rule.

    The rule is a 1.5pt-tall horizontal indigo bar 0.5"-wide that visually
    anchors the tag and separates the slide's header zone from its body.
    """
    tb = add_text(
        slide, text.upper(), left, top, Inches(8), Inches(0.35),
        font=FONT_MONO, size=size, color=color, bold=True,
    )
    if with_rule:
        add_rect(slide, left, top + Inches(0.42),
                 Inches(rule_width), Inches(0.025), color)
    return tb


def add_headline(slide, text, left, top, width, height, *,
                 color=TEXT_PRIMARY, size=FS_H2, align=PP_ALIGN.LEFT,
                 line_spacing=1.08):
    return add_text(
        slide, text, left, top, width, height,
        font=FONT_SANS, size=size, color=color, bold=True,
        line_spacing=line_spacing, align=align,
    )


def add_body(slide, text, left, top, width, height, *,
             color=TEXT_SECONDARY, size=FS_BODY, align=PP_ALIGN.LEFT):
    return add_text(
        slide, text, left, top, width, height,
        font=FONT_SANS, size=size, color=color,
        line_spacing=1.45, align=align,
    )


def add_bullets(slide, bullets, left, top, width, height, *,
                color=TEXT_SECONDARY, size=FS_BODY, marker_color=INDIGO,
                marker="•  "):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.45
        p.space_after = Pt(10)
        m = p.add_run()
        m.text = marker
        m.font.name = FONT_SANS
        m.font.size = Pt(size)
        m.font.color.rgb = marker_color
        m.font.bold = True
        t = p.add_run()
        t.text = bullet
        t.font.name = FONT_SANS
        t.font.size = Pt(size)
        t.font.color.rgb = color
    return tb


def add_footer(slide, citations, *, color=TEXT_MUTED):
    text = "  ·  ".join(citations) if isinstance(citations, list) else citations
    return add_text(
        slide, text, Inches(0.5), Inches(7.05), Inches(12.5), Inches(0.35),
        font=FONT_MONO, size=FS_FOOTER, color=color,
    )


def add_top_corners(
    slide, *, left_tag="POLARIS",
    right_tag="VEEA TRUST TRACK · TECHEX 2026",
    color=TEXT_MUTED,
) -> None:
    add_text(
        slide, left_tag, Inches(0.5), Inches(0.3), Inches(3), Inches(0.3),
        font=FONT_MONO, size=FS_CORNER, color=color, bold=True,
    )
    add_text(
        slide, right_tag, Inches(9.0), Inches(0.3), Inches(4), Inches(0.3),
        font=FONT_MONO, size=FS_CORNER, color=color, align=PP_ALIGN.RIGHT,
    )


def style_table_cell(cell, *, fill, text_color=TEXT_PRIMARY, bold=False,
                     size=FS_BODY_SM, align=PP_ALIGN.LEFT, font=FONT_SANS):
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.margin_left = Inches(0.16)
    cell.margin_right = Inches(0.16)
    cell.margin_top = Inches(0.10)
    cell.margin_bottom = Inches(0.10)
    tf = cell.text_frame
    tf.word_wrap = True
    for p in tf.paragraphs:
        p.alignment = align
        for r in p.runs:
            r.font.name = font
            r.font.size = Pt(size)
            r.font.color.rgb = text_color
            r.font.bold = bold


# --- slide builders -----------------------------------------------------

def slide_1_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    add_top_corners(s)
    add_rect(s, Inches(0), Inches(1.8), Inches(0.08), Inches(4.0), INDIGO)
    add_headline(s, "Polaris", Inches(0.55), Inches(1.7),
                 Inches(8), Inches(1.8), size=FS_HERO, color=TEXT_PRIMARY)
    add_text(s, "PITCH DECK", Inches(0.6), Inches(3.45),
             Inches(8), Inches(0.6),
             font=FONT_MONO, size=16, color=INDIGO, bold=True)
    add_text(s, "From SOC 2 PDF to live AI guardrail in 60 seconds.",
             Inches(0.6), Inches(4.05), Inches(7.5), Inches(0.7),
             font=FONT_SANS, size=FS_LEAD_LARGE, color=TEXT_SECONDARY)
    add_text(s, "MAY 2026  ·  AI & BIG DATA EXPO  ·  SAN JOSE",
             Inches(0.6), Inches(5.0), Inches(8), Inches(0.4),
             font=FONT_MONO, size=FS_LABEL, color=TEXT_MUTED)
    img = DOCS_IMG / "demo_thumbnail.png"
    if img.exists():
        try:
            add_rect(s, Inches(7.55), Inches(1.55), Inches(5.55), Inches(4.4),
                     WHITE, line=BORDER)
            s.shapes.add_picture(str(img), Inches(7.6), Inches(1.6),
                                 width=Inches(5.5))
        except Exception:
            pass
    add_text(s, "Live: polaris--lucaslootan.replit.app",
             Inches(0.6), Inches(6.45), Inches(8), Inches(0.35),
             font=FONT_MONO, size=FS_BODY_SM, color=INDIGO, bold=True)


def slide_2_mission(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, INDIGO_DARK)
    add_text(s, "POLARIS", Inches(0.5), Inches(0.3), Inches(3), Inches(0.3),
             font=FONT_MONO, size=FS_CORNER, color=TEXT_INVERSE_SOFT, bold=True)
    add_text(s, "VEEA TRUST TRACK · TECHEX 2026", Inches(9.0), Inches(0.3),
             Inches(4), Inches(0.3),
             font=FONT_MONO, size=FS_CORNER, color=TEXT_INVERSE_SOFT,
             align=PP_ALIGN.RIGHT)
    add_text(s, "MISSION", Inches(1), Inches(1.5), Inches(11.3), Inches(0.4),
             font=FONT_MONO, size=15, color=TEXT_INVERSE_SOFT, bold=True,
             align=PP_ALIGN.CENTER)
    # Underline rule (centered)
    add_rect(s, Inches((13.333 - 0.6) / 2), Inches(2.0),
             Inches(0.6), Inches(0.03), TEXT_INVERSE_SOFT)
    add_headline(s, "Compile compliance documents\ninto running AI firewalls.",
                 Inches(1), Inches(2.4), Inches(11.3), Inches(2.4),
                 size=FS_H1_LARGE, color=TEXT_INVERSE, align=PP_ALIGN.CENTER)
    add_text(s, "At AI speed.",
             Inches(1), Inches(4.95), Inches(11.3), Inches(0.8),
             font=FONT_SANS, size=38, color=TEXT_INVERSE_SOFT, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(
        s,
        "Polaris is the only end-to-end loop from compliance PDF → deployed runtime policy → continuous adversarial verification.",
        Inches(1.5), Inches(6.05), Inches(10.3), Inches(0.8),
        font=FONT_SANS, size=FS_BODY_LG, color=TEXT_INVERSE_SOFT,
        align=PP_ALIGN.CENTER,
    )


def slide_3_overview(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    add_top_corners(s)
    add_tag(s, "OVERVIEW", Inches(0.6), Inches(0.95))
    add_headline(s, "A 4-agent closed loop that\ncompiles compliance into\nruntime enforcement.",
                 Inches(0.6), Inches(1.65), Inches(7), Inches(3.2),
                 size=FS_H2, color=TEXT_PRIMARY)

    # Highlights panel
    add_rect(s, Inches(7.8), Inches(1.4), Inches(5.0), Inches(2.6),
             WHITE, line=BORDER)
    add_tag(s, "HIGHLIGHTS", Inches(8.0), Inches(1.6),
            color=INDIGO, size=FS_TAG_SM, rule_width=0.35)
    add_bullets(s, [
        "Compiles SOC 2 / HIPAA / EU AI Act / PCI-DSS PDFs into deployable Lobster Trap YAML in ~11s (60s SLA).",
        "Red Team agent continuously stress-tests deployed policy; auto-patches on gap.",
        "Audit-defensible chain of custody mapped to SOC 2 CC8.1.",
    ], Inches(8.0), Inches(2.15), Inches(4.7), Inches(1.95),
       color=TEXT_SECONDARY, size=FS_BODY_XS)

    # Industry & market
    add_rect(s, Inches(7.8), Inches(4.15), Inches(5.0), Inches(1.55),
             WHITE, line=BORDER)
    add_tag(s, "INDUSTRY & MARKET", Inches(8.0), Inches(4.35),
            color=INDIGO, size=FS_TAG_SM, rule_width=0.35)
    add_text(s, "Enterprise AI TRiSM",
             Inches(8.0), Inches(4.9), Inches(4.7), Inches(0.4),
             font=FONT_SANS, size=FS_BODY_SM, color=TEXT_SECONDARY)
    add_text(s, "$2.34B → $7.44B",
             Inches(8.0), Inches(5.2), Inches(4.7), Inches(0.5),
             font=FONT_SANS, size=FS_NUM_LG, color=INDIGO, bold=True)
    add_text(s, "2024 → 2030  ·  21.6% CAGR",
             Inches(8.0), Inches(5.85), Inches(4.7), Inches(0.3),
             font=FONT_MONO, size=FS_LABEL, color=TEXT_MUTED)

    # Team
    add_rect(s, Inches(0.6), Inches(5.0), Inches(7), Inches(1.95),
             WHITE, line=BORDER)
    add_tag(s, "TEAM", Inches(0.8), Inches(5.2),
            color=INDIGO, size=FS_TAG_SM, rule_width=0.35)
    add_text(s, "Lucas Loo Tan Yu Heng — Founder & Lead AI Engineer",
             Inches(0.8), Inches(5.75), Inches(6.7), Inches(0.4),
             font=FONT_SANS, size=FS_BODY, color=TEXT_PRIMARY, bold=True)
    add_body(s, "Day job: Hoppi (M) Sdn Bhd — Hotseller V5 (25+ orchestrated Gemini agents in production).",
             Inches(0.8), Inches(6.15), Inches(6.7), Inches(0.8),
             color=TEXT_SECONDARY, size=FS_BODY_SM)

    # Hero-metric (compact bars on right)
    add_rect(s, Inches(7.8), Inches(5.85), Inches(5.0), Inches(1.1),
             WHITE, line=BORDER)
    add_text(s, "HERO METRIC", Inches(7.95), Inches(5.97),
             Inches(4.7), Inches(0.25),
             font=FONT_MONO, size=FS_TAG_SM, color=INDIGO, bold=True)
    add_text(s, "11s actual", Inches(8.0), Inches(6.27),
             Inches(1.5), Inches(0.25),
             font=FONT_MONO, size=FS_BODY_SM, color=INDIGO, bold=True)
    add_rect(s, Inches(9.6), Inches(6.32), Inches(0.3), Inches(0.18), INDIGO)
    add_text(s, "60s SLA", Inches(8.0), Inches(6.6),
             Inches(1.5), Inches(0.25),
             font=FONT_MONO, size=FS_BODY_SM, color=TEXT_MUTED)
    add_rect(s, Inches(9.6), Inches(6.65), Inches(1.7), Inches(0.18),
             TEXT_MUTED)

    add_footer(s, [
        "Grand View Research · grandviewresearch.com/press-release/global-ai-trust-risk-security-management-market",
    ])


def slide_4_problem_why_now(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    add_top_corners(s)

    # Left panel
    add_rect(s, Inches(0.5), Inches(1.0), Inches(6.1), Inches(5.95),
             WHITE, line=BORDER)
    add_tag(s, "THE PROBLEM", Inches(0.8), Inches(1.3),
            color=DANGER, size=FS_TAG, rule_width=0.55)
    add_headline(s,
                 "Compliance lives in PDFs.\nAI agents run in production.\nThe wire between is hand-written.",
                 Inches(0.8), Inches(2.0), Inches(5.6), Inches(2.4),
                 size=22, color=TEXT_PRIMARY)
    add_bullets(s, [
        "23% of enterprises scaling AI agents; 39% experimenting. (McKinsey State of AI 2025)",
        "233 AI security incidents in 2024 → 362 in 2025, +55% YoY. (Stanford HAI AI Index)",
        "Median compliance / business-formation counsel: $378 / hr. (Clio Legal Trends 2026)",
    ], Inches(0.8), Inches(4.55), Inches(5.6), Inches(2.4),
       size=FS_BODY_SM, color=TEXT_SECONDARY, marker_color=DANGER)

    # Right panel (indigo emphasis)
    add_rect(s, Inches(6.8), Inches(1.0), Inches(6.0), Inches(5.95),
             INDIGO_DARK, line=INDIGO_DARK)
    add_tag(s, "WHY NOW", Inches(7.0), Inches(1.3),
            color=TEXT_INVERSE_SOFT, size=FS_TAG, rule_width=0.55)
    add_headline(s,
                 "Regulators are catching up.\nEvery enterprise running an\nAI agent is non-conformant\nby default.",
                 Inches(7.0), Inches(2.0), Inches(5.6), Inches(2.4),
                 size=22, color=TEXT_INVERSE)
    add_bullets(s, [
        "EU AI Act high-risk: 2 Aug 2026 statutory. Digital Omnibus (7 May 2026) defers Annex III to 2 Dec 2027 — not yet enacted.",
        "Colorado AI Act (SB24-205): postponed to 30 Jun 2026; enforcement stayed pending federal litigation.",
        "NIST AI RMF 1.0 — voluntary framework, Jan 2023.",
        "OWASP LLM Top 10 v1.1 — prompt injection ranked #1.",
    ], Inches(7.0), Inches(4.55), Inches(5.6), Inches(2.4),
       size=FS_BODY_XS, color=TEXT_INVERSE_SOFT, marker_color=TEXT_INVERSE)

    add_footer(s, [
        "mckinsey.com/state-of-ai",
        "hai.stanford.edu/ai-index",
        "clio.com/resources/legal-trends",
        "artificialintelligenceact.eu",
        "leg.colorado.gov/bills/sb24-205",
        "nist.gov · owasp.org",
    ])


def slide_5_solution(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    add_top_corners(s)
    add_tag(s, "SOLUTION", Inches(0.6), Inches(0.95))
    add_headline(s,
                 "Polaris compiles compliance documents\ninto runtime AI guardrails — and verifies them.",
                 Inches(0.6), Inches(1.65), Inches(12.5), Inches(1.5),
                 size=FS_H3, color=TEXT_PRIMARY)
    add_text(s, "Four Gemini agents. One Veea Lobster Trap firewall. One closed loop.",
             Inches(0.6), Inches(3.15), Inches(12.5), Inches(0.4),
             font=FONT_SANS, size=FS_BODY_LG, color=INDIGO, bold=True)

    cards = [
        ("READER AGENT", "gemini-3.1-flash-lite",
         "Extracts compliance requirements from PDF text.",
         "~3s per doc", INDIGO),
        ("SYNTHESIZER", "gemini-3.1-flash-lite + thinking=low",
         "Schema-first synthesis: passes LobsterTrapPolicy Pydantic class as response_schema.",
         "~4.6s median", SUCCESS),
        ("RED TEAM", "gemini-3.1-pro-preview",
         "Generates adversarial probes; triggers Synthesizer regeneration on gap.",
         "~10s / round", WARNING),
    ]
    card_w, card_h = Inches(4.0), Inches(3.4)
    spacing = Inches(0.2)
    total_w = card_w * 3 + spacing * 2
    start_x = (SLIDE_W - total_w) / 2

    for i, (tag, model, body, latency, accent) in enumerate(cards):
        x = start_x + (card_w + spacing) * i
        y = Inches(3.85)
        add_rect(s, x, y, card_w, card_h, WHITE, line=BORDER)
        add_rect(s, x, y, card_w, Inches(0.08), accent)
        add_text(s, tag, x + Inches(0.3), y + Inches(0.3),
                 card_w - Inches(0.6), Inches(0.45),
                 font=FONT_MONO, size=14, color=accent, bold=True)
        add_text(s, model, x + Inches(0.3), y + Inches(0.85),
                 card_w - Inches(0.6), Inches(0.4),
                 font=FONT_MONO, size=FS_BODY_SM, color=TEXT_MUTED)
        add_text(s, body, x + Inches(0.3), y + Inches(1.45),
                 card_w - Inches(0.6), Inches(1.5),
                 font=FONT_SANS, size=FS_BODY_SM, color=TEXT_SECONDARY,
                 line_spacing=1.45)
        add_text(s, latency, x + Inches(0.3), y + card_h - Inches(0.65),
                 card_w - Inches(0.6), Inches(0.4),
                 font=FONT_MONO, size=FS_BODY_LG, color=accent, bold=True)

    add_footer(s, [
        "ai.google.dev/gemini-api/docs/pricing",
        "github.com/veeainc/lobstertrap",
    ])


def slide_6_product(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    add_top_corners(s)
    add_tag(s, "PRODUCT · LIVE DASHBOARD", Inches(0.6), Inches(0.95))
    add_headline(s, "Drag-drop.\n11 seconds.\nLive firewall.",
                 Inches(0.6), Inches(1.65), Inches(6.5), Inches(3.0),
                 size=FS_H1, color=TEXT_PRIMARY)
    add_body(s,
             "Phase-9 schema-first Synthesizer cuts latency 5× while tying with the larger Pro model on intrinsic accuracy (6.0/11 LT corpus). Same architectural pattern shipped at scale in Hoppi's Hotseller V5 (25+ orchestrated agents).",
             Inches(0.6), Inches(4.7), Inches(6.3), Inches(2.0),
             color=TEXT_SECONDARY, size=FS_BODY_LG)

    img = DOCS_IMG / "demo_thumbnail.png"
    if img.exists():
        try:
            add_rect(s, Inches(7.05), Inches(1.55), Inches(5.85), Inches(4.4),
                     WHITE, line=BORDER)
            s.shapes.add_picture(str(img), Inches(7.15), Inches(1.65),
                                 width=Inches(5.65))
        except Exception:
            pass

    add_footer(s, [
        "docs/MODEL_BAKEOFF.md (48-run model bake-off, Phase 9, 2026-05-13)",
        "github.com/seekerPrice/polaris",
    ])


def slide_7_closed_loop(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    add_top_corners(s)
    add_tag(s, "THE CLOSED LOOP", Inches(0.6), Inches(0.95))
    add_headline(s, "This loop closes itself.",
                 Inches(0.6), Inches(1.65), Inches(12), Inches(0.9),
                 size=FS_H2_LARGE, color=TEXT_PRIMARY)
    add_text(s, "Lobster Trap _lobstertrap declared-vs-detected mismatches feed the Red Team. Humans on the audit trail.",
             Inches(0.6), Inches(2.65), Inches(12), Inches(0.4),
             font=FONT_SANS, size=FS_LEAD_SM, color=TEXT_SECONDARY)

    steps = [
        ("1", "Reader", "~3s", INDIGO),
        ("2", "Synthesizer", "4.6s · schema-first", SUCCESS),
        ("3", "Validation", "11/11 LT corpus", VIOLET),
        ("4", "Lobster Trap", "_lobstertrap intent", WARNING),
        ("5", "Red Team", "~10s · auto-patch", DANGER),
    ]
    box_w, box_h = Inches(2.2), Inches(2.4)
    arrow_w = Inches(0.18)
    total_w = box_w * 5 + arrow_w * 4
    start_x = (SLIDE_W - total_w) / 2

    for i, (n, label, sub, accent) in enumerate(steps):
        x = start_x + (box_w + arrow_w) * i
        y = Inches(3.4)
        add_rect(s, x, y, box_w, box_h, WHITE, line=BORDER)
        add_rect(s, x, y, box_w, Inches(0.08), accent)
        add_text(s, n, x + Inches(0.2), y + Inches(0.3),
                 Inches(1.5), Inches(0.7),
                 font=FONT_MONO, size=FS_ICON_LG, color=accent, bold=True)
        add_text(s, label, x + Inches(0.2), y + Inches(1.05),
                 box_w - Inches(0.4), Inches(0.5),
                 font=FONT_SANS, size=15, color=TEXT_PRIMARY, bold=True)
        add_text(s, sub, x + Inches(0.2), y + Inches(1.6),
                 box_w - Inches(0.4), Inches(0.5),
                 font=FONT_MONO, size=11, color=TEXT_SECONDARY)
        if i < 4:
            ax = x + box_w
            add_text(s, "→", ax, y + Inches(1.05),
                     arrow_w, Inches(0.5),
                     font=FONT_MONO, size=20, color=TEXT_MUTED, bold=True,
                     align=PP_ALIGN.CENTER)

    add_text(s,
             "↻  Synthesizer regenerates on gap → hot-reload → same probe now blocked.",
             Inches(1), Inches(6.2), Inches(11.3), Inches(0.5),
             font=FONT_MONO, size=14, color=INDIGO, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(s, "AI governing AI — with humans on the audit trail.",
             Inches(1), Inches(6.65), Inches(11.3), Inches(0.4),
             font=FONT_SANS, size=FS_BODY_SM, color=TEXT_MUTED,
             align=PP_ALIGN.CENTER)


def slide_8_market(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    add_top_corners(s)
    add_tag(s, "MARKET OPPORTUNITY", Inches(0.6), Inches(0.95))
    add_headline(s, "A $7.44B market,\nexpanding 21.6% per year.",
                 Inches(0.6), Inches(1.65), Inches(7), Inches(1.9),
                 size=FS_H2, color=TEXT_PRIMARY)

    bullets = [
        ("TAM", "$7.44B by 2030",
         "Global AI TRiSM. $2.34B (2024) → $7.44B at 21.6% CAGR. (Grand View Research, 2025)"),
        ("SAM", "Billion-dollar (Gartner Feb 2026)",
         "AI governance platforms. Precedence: $309M → $3.59B by 2033 at 36% CAGR."),
        ("TARGET", "$740M serviceable",
         "US SOC 2 / HIPAA enterprises with AI agents in production (822k HIPAA × 23% McKinsey AI-in-prod rate)."),
        ("ENTRY", "$74M ARR potential",
         "1% capture × $500–$2,499 / mo / policy ARPU."),
    ]
    y = Inches(3.55)
    for tag, big, body in bullets:
        add_text(s, tag, Inches(0.6), y,
                 Inches(1.1), Inches(0.4),
                 font=FONT_MONO, size=FS_TAG, color=INDIGO, bold=True)
        add_text(s, big, Inches(1.75), y - Inches(0.06),
                 Inches(5.5), Inches(0.5),
                 font=FONT_SANS, size=FS_NUM_MD, color=TEXT_PRIMARY, bold=True)
        add_text(s, body, Inches(0.6), y + Inches(0.48),
                 Inches(6.8), Inches(0.5),
                 font=FONT_SANS, size=FS_BODY_SM, color=TEXT_SECONDARY,
                 line_spacing=1.35)
        y += Inches(0.85)

    # Concentric circles right
    cx = Inches(10.0)
    cy = Inches(4.5)
    sizes = [
        (5.4, INDIGO_SOFT, BORDER, "TAM\n$7.44B", TEXT_PRIMARY, 0.5, 12),
        (3.9, RGBColor(0xC7, 0xD2, 0xFE), INDIGO_SOFT, "SAM\n$3.6B",
         TEXT_PRIMARY, 0.9, 12),
        (2.5, INDIGO_LIGHT, INDIGO, "TARGET\n$740M",
         TEXT_INVERSE, 1.45, 12),
        (1.2, INDIGO_DARK, INDIGO_DARK, "$74M",
         TEXT_INVERSE, 2.05, 14),
    ]
    for d, fill, line, label, text_color, ty_inch, font_pt in sizes:
        diam = Inches(d)
        add_oval(s, cx - diam / 2, cy - diam / 2, diam, diam, fill,
                 line=line, line_width=Pt(0.75))
        add_text(s, label, cx - Inches(1.5),
                 cy - diam / 2 + Inches(ty_inch),
                 Inches(3), Inches(0.7),
                 font=FONT_MONO, size=font_pt, color=text_color, bold=True,
                 align=PP_ALIGN.CENTER, line_spacing=1.15)

    add_footer(s, [
        "grandviewresearch.com · gartner.com (Feb 2026 PR)",
        "precedenceresearch.com",
        "mckinsey.com/state-of-ai",
        "hhs.gov/hipaa",
    ])


def slide_9_competitive(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    add_top_corners(s)
    add_tag(s, "COMPETITIVE LANDSCAPE", Inches(0.6), Inches(0.95))
    add_headline(s, "Polaris is the only end-to-end\nPDF → Deploy → Verify → Patch loop.",
                 Inches(0.6), Inches(1.65), Inches(12), Inches(1.4),
                 size=FS_H3, color=TEXT_PRIMARY)

    headers = ["Capability", "Polaris", "MS Purview\nAI Hub", "Lakera\nGuard",
               "F5 AI Guardrails\n(ex-CalypsoAI)", "Cisco AI Defense\n(ex-Robust Intelligence)"]
    rows = [
        ("Auto-synthesize policy from compliance PDF",
         "✓", "✗", "✗", "✗", "✗"),
        ("Runtime inline enforcement on LLM I/O",
         "✓", "✗*", "✓", "✓", "✓"),
        ("Closed-loop continuous Red Team",
         "✓", "✗", "✗", "✗", "✗"),
        ("Auto-generated compliance PDF",
         "✓", "partial", "✗", "✗", "partial"),
        ("Open-source DPI substrate",
         "✓ (Veea LT)", "✗", "✗", "✗", "✗"),
    ]
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = s.shapes.add_table(n_rows, n_cols,
                                     Inches(0.6), Inches(3.4),
                                     Inches(12.1), Inches(3.0))
    table = table_shape.table
    col_widths = [3.6, 1.6, 1.4, 1.4, 2.0, 2.1]
    for i, w in enumerate(col_widths):
        table.columns[i].width = Inches(w)

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        is_polaris_col = (i == 1)
        style_table_cell(
            cell, fill=INDIGO if is_polaris_col else INDIGO_DARK,
            text_color=TEXT_INVERSE, bold=True, size=FS_BODY_SM,
            align=PP_ALIGN.CENTER, font=FONT_MONO,
        )

    for r_idx, row in enumerate(rows, start=1):
        alt = (r_idx % 2 == 0)
        row_bg = SUBTLE if alt else WHITE
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = val
            is_polaris_col = (c_idx == 1)
            cell_bg = INDIGO_SOFT if is_polaris_col else row_bg
            if c_idx == 0:
                style_table_cell(cell, fill=cell_bg, text_color=TEXT_PRIMARY,
                                 size=FS_BODY_SM, align=PP_ALIGN.LEFT)
            else:
                if val == "✓":
                    style_table_cell(cell, fill=cell_bg, text_color=SUCCESS,
                                     bold=True, size=FS_ICON,
                                     align=PP_ALIGN.CENTER, font=FONT_SANS)
                elif val.startswith("✗"):
                    style_table_cell(cell, fill=cell_bg, text_color=DANGER,
                                     bold=True, size=FS_ICON,
                                     align=PP_ALIGN.CENTER, font=FONT_SANS)
                elif val == "partial":
                    style_table_cell(cell, fill=cell_bg, text_color=WARNING,
                                     bold=True, size=FS_BODY_SM,
                                     align=PP_ALIGN.CENTER, font=FONT_MONO)
                else:
                    style_table_cell(cell, fill=cell_bg, text_color=SUCCESS,
                                     bold=True, size=FS_BODY_SM,
                                     align=PP_ALIGN.CENTER, font=FONT_MONO)

    add_text(s, "* Microsoft Purview AI Hub: endpoint browser DLP + post-hoc audit only — not inline blocking on LLM outputs.",
             Inches(0.6), Inches(6.5), Inches(12.5), Inches(0.3),
             font=FONT_SANS, size=FS_LABEL, color=TEXT_MUTED)
    add_text(s, "None of the four funded incumbents combine auto-synthesis + runtime enforcement + closed-loop verification.",
             Inches(0.6), Inches(6.82), Inches(12.5), Inches(0.3),
             font=FONT_SANS, size=FS_BODY_SM, color=INDIGO, bold=True)

    add_footer(s, [
        "learn.microsoft.com/purview",
        "lakera.ai/lakera-guard",
        "f5.com/products/ai-guardrails",
        "cisco.com/products/security/ai-defense",
    ])


def slide_10_business_model(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    add_top_corners(s)
    add_tag(s, "BUSINESS MODEL", Inches(0.6), Inches(0.95))
    add_headline(s, "Per-policy subscription.\nSOC 2-defensible.\n>95% gross margin.",
                 Inches(0.6), Inches(1.65), Inches(6.5), Inches(3.5),
                 size=FS_H3, color=TEXT_PRIMARY)
    add_body(s,
             "Tiered SaaS + per-inference overage. Self-hosted Lobster Trap option for compliance-strict tenants. Pricing is a draft model pending design-partner validation (Q3 2026 pilots).",
             Inches(0.6), Inches(5.4), Inches(6.3), Inches(1.6),
             color=TEXT_SECONDARY, size=FS_BODY_SM)

    headers = ["", "Starter", "Pro"]
    rows = [
        ("Compliance packs", "1", "Unlimited"),
        ("AI agents protected", "1", "Unlimited"),
        ("Continuous Red Team", "weekly", "continuous"),
        ("Drift detection (v0.2)", "—", "✓"),
        ("Slack / Linear alerts", "—", "✓"),
        ("Compliance PDF auto-gen", "✓", "✓"),
        ("Self-hosted LT option", "—", "✓"),
        ("Audience", "SOC 2 design partners", "Multi-policy enterprises"),
        ("Price", "$499 / mo", "$2,499 / mo + $0.10/1k"),
    ]
    n_rows = len(rows) + 1

    table_shape = s.shapes.add_table(n_rows, 3,
                                     Inches(7.3), Inches(1.5),
                                     Inches(5.5), Inches(5.2))
    table = table_shape.table
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(1.5)
    table.columns[2].width = Inches(1.8)

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        is_pro = (i == 2)
        style_table_cell(
            cell, fill=INDIGO if is_pro else (INDIGO_DARK if i == 1 else SUBTLE),
            text_color=TEXT_INVERSE if i > 0 else TEXT_MUTED,
            bold=True, size=15, align=PP_ALIGN.CENTER, font=FONT_MONO,
        )

    for r_idx, row in enumerate(rows, start=1):
        is_price = (row[0] == "Price")
        alt = (r_idx % 2 == 0)
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = val
            row_bg = SUBTLE if alt else WHITE
            if c_idx == 0:
                style_table_cell(cell, fill=row_bg,
                                 text_color=TEXT_PRIMARY if is_price else TEXT_SECONDARY,
                                 bold=is_price, size=FS_BODY_SM)
            else:
                if is_price:
                    text_color = INDIGO
                elif val == "—":
                    text_color = TEXT_MUTED
                elif val == "✓":
                    text_color = SUCCESS
                else:
                    text_color = TEXT_SECONDARY
                style_table_cell(cell, fill=row_bg, text_color=text_color,
                                 bold=is_price or val == "✓",
                                 size=FS_BODY if not is_price else 15,
                                 align=PP_ALIGN.CENTER, font=FONT_MONO)

    add_text(s,
             "Marginal cost / policy ≈ $0.005 in Gemini compute (gemini-3.1-flash-lite). Gross margin >95%.",
             Inches(0.6), Inches(6.8), Inches(6.5), Inches(0.4),
             font=FONT_MONO, size=FS_BODY_SM, color=INDIGO, bold=True)

    add_footer(s, ["ai.google.dev/gemini-api/docs/pricing"])


def slide_11_unit_economics(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    add_top_corners(s)
    add_tag(s, "UNIT ECONOMICS", Inches(0.6), Inches(0.95))
    add_headline(s, "Why the math works.",
                 Inches(0.6), Inches(1.65), Inches(12), Inches(0.9),
                 size=FS_H2_LARGE, color=TEXT_PRIMARY)

    tiles = [
        ("11 seconds", "END-TO-END  ·  60s SLA", INDIGO),
        ("$0.005", "GEMINI COST / POLICY", SUCCESS),
        ("3,000,000×", "COST COMPRESSION", WARNING),
        ("$15K–$45K", "COMPLIANCE COUNSEL REPLACED", VIOLET),
        ("11 / 11", "LOBSTER TRAP CORPUS PASS", INFO),
        ("62 / 62", "UNIT TESTS PASS", SUCCESS),
    ]
    cols = 3
    tile_w, tile_h = Inches(4.0), Inches(2.05)
    spacing_x = Inches(0.15)
    spacing_y = Inches(0.2)
    total_w = tile_w * cols + spacing_x * (cols - 1)
    start_x = (SLIDE_W - total_w) / 2
    start_y = Inches(2.85)

    for i, (big, cap, accent) in enumerate(tiles):
        r, c = divmod(i, cols)
        x = start_x + (tile_w + spacing_x) * c
        y = start_y + (tile_h + spacing_y) * r
        add_rect(s, x, y, tile_w, tile_h, WHITE, line=BORDER)
        add_rect(s, x, y, Inches(0.1), tile_h, accent)
        add_text(s, big, x + Inches(0.4), y + Inches(0.35),
                 tile_w - Inches(0.6), Inches(0.95),
                 font=FONT_SANS, size=FS_NUM_XL, color=TEXT_PRIMARY, bold=True)
        add_text(s, cap, x + Inches(0.4), y + Inches(1.4),
                 tile_w - Inches(0.6), Inches(0.5),
                 font=FONT_MONO, size=FS_LABEL_BOLD, color=accent, bold=True)

    add_text(s,
             "Counsel-cost floor: Drata SOC 2 framework (3-wk observation + 2-5wk audit) × Clio's $378/hr 2026 median compliance rate × FTE ≈ $45K per policy cycle.",
             Inches(0.6), Inches(6.75), Inches(12.1), Inches(0.4),
             font=FONT_SANS, size=FS_LABEL, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    add_footer(s, [
        "drata.com/grc-central/soc-2",
        "clio.com/resources/legal-trends",
        "ai.google.dev/gemini-api/docs/pricing",
        "github.com/seekerPrice/polaris",
    ])


def slide_12_roadmap(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    add_top_corners(s)
    add_tag(s, "ROADMAP", Inches(0.6), Inches(0.95))
    add_headline(s, "Where we go after the hackathon.",
                 Inches(0.6), Inches(1.65), Inches(12), Inches(0.9),
                 size=FS_H3, color=TEXT_PRIMARY)

    quarters = [
        ("Q3 2026", "next · pilots", [
            "Per-agent declared_intent verdicts (v0.2)",
            "SOC 2 design-partner pilot (3-5 enterprises)",
            "Drift detection v0.1",
        ], INDIGO),
        ("Q4 2026", "GA", [
            "HIPAA + PCI-DSS pack hardening",
            "Pricing GA (Starter + Pro)",
            "Slack / Linear alerts",
        ], INDIGO_LIGHT),
        ("Q1 2027", "SEED · $2M", [
            "Multi-tenant SaaS architecture",
            "$2M Seed target",
        ], WARNING),
        ("Q2 2027", "scale", [
            "EU AI Act Annex III templates",
            "25 paying tenants",
        ], INDIGO_LIGHT),
        ("Q3 2027", "channel", [
            "Self-service Veea DevKit deploy",
            "100+ tenants",
        ], INDIGO),
    ]
    col_w = Inches(2.4)
    gap = Inches(0.05)
    total_w = col_w * 5 + gap * 4
    start_x = (SLIDE_W - total_w) / 2
    col_y = Inches(2.9)
    col_h = Inches(3.65)

    for i, (q, sub, items, accent) in enumerate(quarters):
        x = start_x + (col_w + gap) * i
        add_rect(s, x, col_y, col_w, col_h, WHITE, line=BORDER)
        add_rect(s, x, col_y, col_w, Inches(0.08), accent)
        add_text(s, q, x + Inches(0.2), col_y + Inches(0.28),
                 col_w - Inches(0.4), Inches(0.5),
                 font=FONT_SANS, size=20, color=TEXT_PRIMARY, bold=True)
        add_text(s, sub, x + Inches(0.2), col_y + Inches(0.85),
                 col_w - Inches(0.4), Inches(0.4),
                 font=FONT_MONO, size=FS_LABEL, color=accent, bold=True)
        add_bullets(s, items,
                    x + Inches(0.2), col_y + Inches(1.35),
                    col_w - Inches(0.4), col_h - Inches(1.5),
                    color=TEXT_SECONDARY, size=FS_LABEL, marker_color=accent)

    bar_x = start_x + (col_w + gap) * 2
    bar_y = col_y + col_h + Inches(0.15)
    add_rect(s, bar_x, bar_y, col_w, Inches(0.2), WARNING)
    add_text(s, "FUNDRAISING", bar_x, bar_y + Inches(0.24),
             col_w, Inches(0.3),
             font=FONT_MONO, size=FS_LABEL_BOLD, color=WARNING, bold=True,
             align=PP_ALIGN.CENTER)


def slide_13_team(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    add_top_corners(s)
    add_tag(s, "TEAM", Inches(0.6), Inches(0.95))
    add_headline(s, "Sole engineer, with the\nproduction playbook already shipped.",
                 Inches(0.6), Inches(1.65), Inches(12), Inches(1.4),
                 size=FS_H3, color=TEXT_PRIMARY)

    add_rect(s, Inches(0.6), Inches(3.3), Inches(7.5), Inches(3.55),
             WHITE, line=BORDER)
    add_rect(s, Inches(0.6), Inches(3.3), Inches(0.1), Inches(3.55), INDIGO)
    add_oval(s, Inches(0.9), Inches(3.55), Inches(1.6), Inches(1.6),
             INDIGO, line=INDIGO, line_width=Pt(1))
    add_text(s, "LT", Inches(0.9), Inches(3.85), Inches(1.6), Inches(1.0),
             font=FONT_SANS, size=44, color=TEXT_INVERSE, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(s, "Lucas — Loo Tan Yu Heng",
             Inches(2.7), Inches(3.55), Inches(5.3), Inches(0.5),
             font=FONT_SANS, size=FS_H4, color=TEXT_PRIMARY, bold=True)
    add_text(s, "FOUNDER & LEAD AI ENGINEER  ·  KUALA LUMPUR, MY",
             Inches(2.7), Inches(4.12), Inches(5.3), Inches(0.4),
             font=FONT_MONO, size=FS_LABEL, color=INDIGO, bold=True)
    add_body(s,
             "Lead AI Engineer at Hoppi (M) Sdn Bhd on Hotseller V5 — 25+ orchestrated Gemini agents, multi-tier model routing, semantic cache invalidation, taxonomy classifier across 51 categories × 100K+ multilingual records. Authored 26-entry LLM Production Anti-Pattern Registry. Polaris's 4-agent loop is the same pattern, applied to compliance.",
             Inches(2.7), Inches(4.65), Inches(5.3), Inches(2.1),
             color=TEXT_SECONDARY, size=FS_BODY_SM)

    advisors = [
        ("CISO ADVISOR", "OPEN"),
        ("COMPLIANCE COUNSEL", "OPEN"),
        ("VEEA PARTNERSHIPS", "OPEN"),
    ]
    card_w = Inches(1.55)
    card_h = Inches(3.55)
    card_y = Inches(3.3)
    start_x = Inches(8.4)
    gap = Inches(0.2)
    for i, (role, status) in enumerate(advisors):
        x = start_x + (card_w + gap) * i
        add_rect(s, x, card_y, card_w, card_h, WHITE, line=BORDER)
        add_oval(s, x + Inches(0.275), card_y + Inches(0.4),
                 Inches(1.0), Inches(1.0), SUBTLE,
                 line=BORDER_STRONG, line_width=Pt(0.75))
        add_text(s, "?", x + Inches(0.275), card_y + Inches(0.55),
                 Inches(1.0), Inches(0.8),
                 font=FONT_SANS, size=36, color=TEXT_MUTED, bold=True,
                 align=PP_ALIGN.CENTER)
        add_text(s, role, x + Inches(0.15), card_y + Inches(1.7),
                 card_w - Inches(0.3), Inches(0.8),
                 font=FONT_MONO, size=FS_LABEL_BOLD, color=TEXT_PRIMARY,
                 bold=True, align=PP_ALIGN.CENTER, line_spacing=1.2)
        add_text(s, status, x + Inches(0.15), card_y + Inches(2.55),
                 card_w - Inches(0.3), Inches(0.4),
                 font=FONT_MONO, size=FS_BODY_SM, color=WARNING, bold=True,
                 align=PP_ALIGN.CENTER)
        add_text(s, "recruiting", x + Inches(0.15), card_y + Inches(3.0),
                 card_w - Inches(0.3), Inches(0.4),
                 font=FONT_MONO, size=FS_LABEL, color=TEXT_MUTED,
                 align=PP_ALIGN.CENTER)

    add_footer(s, ["github.com/seekerPrice/polaris", "linkedin.com/in/lucasloo"])


def slide_14_ask(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    add_top_corners(s)
    add_tag(s, "WHAT WE'RE ASKING", Inches(0.6), Inches(0.95))
    add_headline(s,
                 "Hackathon-first.\nPilot deployments next.\n$2M Seed in Q1 2027.",
                 Inches(0.6), Inches(1.65), Inches(7), Inches(3.0),
                 size=FS_H2, color=TEXT_PRIMARY)
    add_body(s,
             "Today's ask is operational, not financial: deploy Polaris as a Veea DevKit pilot in your enterprise. Goal: 3–5 SOC 2 design partners in Q3 2026.",
             Inches(0.6), Inches(4.7), Inches(6.5), Inches(1.5),
             color=TEXT_SECONDARY, size=FS_BODY_LG)

    bars = [
        ("ENGINEERING (3 FTE)", 50, INDIGO),
        ("SALES & PARTNERSHIPS", 30, SUCCESS),
        ("R&D · drift + multi-tenant", 20, WARNING),
    ]
    add_text(s, "Future Seed allocation ($2M)",
             Inches(0.6), Inches(6.2), Inches(7), Inches(0.4),
             font=FONT_MONO, size=FS_LABEL_BOLD, color=TEXT_MUTED, bold=True)
    bar_total_w = Inches(6.5)
    bar_h = Inches(0.5)
    bar_y = Inches(6.6)
    bar_x = Inches(0.6)
    for label, pct, accent in bars:
        w = bar_total_w * (pct / 100)
        add_rect(s, bar_x, bar_y, w, bar_h, accent)
        add_text(s, f"{pct}%", bar_x + w / 2 - Inches(0.3),
                 bar_y + Inches(0.1),
                 Inches(0.6), Inches(0.3),
                 font=FONT_MONO, size=14, color=TEXT_INVERSE, bold=True,
                 align=PP_ALIGN.CENTER)
        bar_x += w

    add_rect(s, Inches(8.3), Inches(1.5), Inches(4.5), Inches(5.0),
             INDIGO_SOFT, line=INDIGO_LIGHT)
    add_rect(s, Inches(8.3), Inches(1.5), Inches(4.5), Inches(0.1), INDIGO)
    add_tag(s, "IDEAL PARTNER", Inches(8.5), Inches(1.85),
            color=INDIGO_DARK, size=FS_TAG)
    add_bullets(s, [
        "A Veea ecosystem deployment lead (Lobster Trap is the substrate).",
        "A regulated-industry early-adopter (Fortune 2000 in healthcare or fintech).",
        "A counsel-side advisor on EU AI Act high-risk classification (Annex III).",
    ], Inches(8.5), Inches(2.65), Inches(4.1), Inches(3.8), size=FS_BODY,
       color=TEXT_PRIMARY, marker_color=INDIGO)


def slide_15_thank_you(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, INDIGO_DARK)
    add_text(s, "POLARIS", Inches(0.5), Inches(0.3), Inches(3), Inches(0.4),
             font=FONT_MONO, size=FS_CORNER, color=TEXT_INVERSE_SOFT, bold=True)
    add_text(s, "VEEA TRUST TRACK  ·  TECHEX 2026", Inches(9.0), Inches(0.3),
             Inches(4), Inches(0.4),
             font=FONT_MONO, size=FS_CORNER, color=TEXT_INVERSE_SOFT,
             align=PP_ALIGN.RIGHT)
    add_headline(s, "Thank you.",
                 Inches(0.5), Inches(2.4), Inches(12.33), Inches(2.4),
                 size=FS_DISPLAY, color=TEXT_INVERSE, align=PP_ALIGN.CENTER)
    add_text(s, "polaris--lucaslootan.replit.app",
             Inches(0.5), Inches(5.0), Inches(12.33), Inches(0.6),
             font=FONT_MONO, size=24, color=TEXT_INVERSE, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(s, "lucaslootan@gmail.com  ·  github.com/seekerPrice/polaris",
             Inches(0.5), Inches(5.8), Inches(12.33), Inches(0.5),
             font=FONT_MONO, size=FS_BODY_LG, color=TEXT_INVERSE_SOFT,
             align=PP_ALIGN.CENTER)
    add_text(s, "Drop your SOC 2 PDF. Watch the firewall deploy in 11 seconds.",
             Inches(0.5), Inches(6.4), Inches(12.33), Inches(0.5),
             font=FONT_SANS, size=FS_BODY_LG, color=TEXT_INVERSE_SOFT,
             align=PP_ALIGN.CENTER)


def slide_a1_qr(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    add_top_corners(s)
    add_tag(s, "APPENDIX  ·  LIVE DEMO", Inches(0.6), Inches(0.95))
    add_headline(s, "Try Polaris yourself.",
                 Inches(0.6), Inches(1.65), Inches(12), Inches(1.0),
                 size=FS_H2, color=TEXT_PRIMARY)
    add_text(s, "Drop your SOC 2 PDF. Watch the firewall deploy in 11 seconds.",
             Inches(0.6), Inches(2.7), Inches(12), Inches(0.5),
             font=FONT_SANS, size=FS_LEAD, color=TEXT_SECONDARY)

    qr_path = DOCS_IMG / "polaris_qr.png"
    qr_size = Inches(3.5)
    qr_x = (SLIDE_W - qr_size) / 2
    qr_y = Inches(3.5)
    add_rect(s, qr_x - Inches(0.25), qr_y - Inches(0.25),
             qr_size + Inches(0.5), qr_size + Inches(0.5),
             WHITE, line=INDIGO, line_width=Pt(1.5))
    if qr_path.exists():
        try:
            s.shapes.add_picture(str(qr_path), qr_x, qr_y, width=qr_size)
        except Exception:
            pass

    add_text(s, "polaris--lucaslootan.replit.app",
             Inches(0.5), Inches(6.95), Inches(12.33), Inches(0.45),
             font=FONT_MONO, size=22, color=INDIGO, bold=True,
             align=PP_ALIGN.CENTER)


def slide_a2_architecture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    add_top_corners(s)
    add_tag(s, "APPENDIX  ·  ARCHITECTURE", Inches(0.6), Inches(0.95))
    add_headline(s, "The 4-agent + Lobster Trap loop.",
                 Inches(0.6), Inches(1.65), Inches(12), Inches(1.0),
                 size=FS_H3, color=TEXT_PRIMARY)

    # ── Top row (forward flow): PDF → Reader → Synth → policy → LT → Demo Agent
    top_boxes = [
        ("Compliance\nPDF", 0.6, 3.1, 1.7, 1.3, TEXT_MUTED),
        ("READER\nAGENT", 2.7, 3.1, 1.7, 1.3, INDIGO),
        ("SYNTHESIZER\nAGENT", 4.8, 3.1, 1.7, 1.3, SUCCESS),
        ("policy.yaml\n+ schemas", 6.9, 3.1, 1.7, 1.3, VIOLET),
        ("LOBSTER\nTRAP DPI", 9.0, 3.1, 1.7, 1.3, WARNING),
        ("Demo\nAgent", 11.1, 3.1, 1.7, 1.3, TEXT_MUTED),
    ]
    # ── Bottom row (loopback): Mismatch Detector → Red Team
    bottom_boxes = [
        ("MISMATCH\nDETECTOR", 5.85, 5.1, 1.7, 1.05, VIOLET),
        ("RED TEAM\nAGENT", 8.0, 5.1, 1.7, 1.05, DANGER),
    ]

    for label, x, y, w, h, accent in top_boxes + bottom_boxes:
        add_rect(s, Inches(x), Inches(y), Inches(w), Inches(h),
                 WHITE, line=accent, line_width=Pt(1))
        add_text(s, label, Inches(x), Inches(y) + Inches(0.2),
                 Inches(w), Inches(h),
                 font=FONT_MONO, size=FS_LABEL, color=accent, bold=True,
                 align=PP_ALIGN.CENTER, line_spacing=1.2)

    # ── Forward arrows between top-row boxes (filled RIGHT_ARROW shapes)
    for gap_x in [2.3, 4.4, 6.5, 8.6, 10.7]:
        arr = s.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(gap_x + 0.02), Inches(3.55),
            Inches(0.36), Inches(0.4),
        )
        arr.fill.solid()
        arr.fill.fore_color.rgb = INDIGO
        arr.line.fill.background()

    # ── Loop-back arrows (real connector lines with arrowheads)
    # 1. LOBSTER TRAP (bottom center 9.85, 4.4) → MISMATCH DETECTOR (top center 6.7, 5.1)
    add_arrow(
        s,
        begin=(Inches(9.85), Inches(4.4)),
        end=(Inches(6.7), Inches(5.1)),
        color=WARNING, width=Pt(2),
        connector_type=MSO_CONNECTOR.ELBOW,
    )
    add_text(s, "audit log",
             Inches(7.4), Inches(4.5),
             Inches(2), Inches(0.3),
             font=FONT_MONO, size=11, color=WARNING, bold=True,
             align=PP_ALIGN.CENTER)

    # 2. MISMATCH DETECTOR (right edge 7.55, 5.625) → RED TEAM (left edge 8.0, 5.625)
    add_arrow(
        s,
        begin=(Inches(7.55), Inches(5.625)),
        end=(Inches(8.0), Inches(5.625)),
        color=VIOLET, width=Pt(2),
    )

    # 3. RED TEAM (top center 8.85, 5.1) → SYNTHESIZER (bottom center 5.65, 4.4)
    add_arrow(
        s,
        begin=(Inches(8.85), Inches(5.1)),
        end=(Inches(5.65), Inches(4.4)),
        color=DANGER, width=Pt(2),
        connector_type=MSO_CONNECTOR.ELBOW,
    )
    add_text(s, "regenerate on gap",
             Inches(4.5), Inches(4.5),
             Inches(2.5), Inches(0.3),
             font=FONT_MONO, size=11, color=DANGER, bold=True,
             align=PP_ALIGN.CENTER)

    # ── Caption
    add_text(s,
             "Mismatch Detector compares Lobster Trap's _lobstertrap declared-intent vs detected. Gaps fire the Red Team probe, which triggers Synthesizer regeneration. Full architecture in CLAUDE.md §3.",
             Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.6),
             font=FONT_SANS, size=FS_BODY_SM, color=TEXT_MUTED,
             align=PP_ALIGN.CENTER)

    add_footer(s, [
        "github.com/seekerPrice/polaris",
        "github.com/veeainc/lobstertrap",
    ])


def slide_a3_compliance(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    add_top_corners(s)
    add_tag(s, "APPENDIX  ·  COMPLIANCE COVERAGE", Inches(0.6), Inches(0.95))
    add_headline(s, "Every rule traces to a named, citable control.",
                 Inches(0.6), Inches(1.65), Inches(12), Inches(1.0),
                 size=FS_H4, color=TEXT_PRIMARY)

    headers = ["Source control", "Polaris rule examples", "LT action"]
    rows = [
        ("SOC 2 CC6.1 (Logical access)",
         "block_credential_exfiltration", "DENY"),
        ("SOC 2 CC8.1 (Change management)",
         "ApprovalGate consent before deploy", "HUMAN_REVIEW"),
        ("HIPAA §164.312(a)(2) (Access control)",
         "block_phi_unauthorized", "DENY"),
        ("EU AI Act Art. 9 (Risk management)",
         "quarantine_borderline_credential", "QUARANTINE"),
        ("OWASP LLM01 (Prompt Injection)",
         "block_obfuscation_attempts", "DENY"),
        ("OWASP LLM06 (Sensitive Info Disclosure)",
         "egress DPI scan rules", "LOG + DENY"),
    ]
    n_rows = len(rows) + 1
    table_shape = s.shapes.add_table(n_rows, 3,
                                     Inches(0.6), Inches(2.95),
                                     Inches(12.1), Inches(3.6))
    table = table_shape.table
    table.columns[0].width = Inches(4.5)
    table.columns[1].width = Inches(5.0)
    table.columns[2].width = Inches(2.6)
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        style_table_cell(cell, fill=INDIGO_DARK, text_color=TEXT_INVERSE,
                         bold=True, size=FS_BODY, align=PP_ALIGN.LEFT,
                         font=FONT_MONO)
    for r_idx, row in enumerate(rows, start=1):
        alt = (r_idx % 2 == 0)
        row_bg = SUBTLE if alt else WHITE
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = val
            if c_idx == 0:
                style_table_cell(cell, fill=row_bg, text_color=TEXT_SECONDARY,
                                 size=FS_BODY_SM)
            elif c_idx == 1:
                style_table_cell(cell, fill=row_bg, text_color=TEXT_PRIMARY,
                                 size=FS_BODY_SM, font=FONT_MONO)
            else:
                color = DANGER if "DENY" in val else (
                    WARNING if "HUMAN" in val else (
                        VIOLET if "QUARANTINE" in val else SUCCESS))
                style_table_cell(cell, fill=row_bg, text_color=color,
                                 bold=True, size=FS_BODY_SM, font=FONT_MONO,
                                 align=PP_ALIGN.CENTER)

    add_footer(s, [
        "aicpa-cima.com (SOC 2 TSC)",
        "hhs.gov/hipaa · 45 CFR §164.312",
        "artificialintelligenceact.eu",
        "owasp.org",
    ])


# --- main ---------------------------------------------------------------

def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_1_cover(prs)
    slide_2_mission(prs)
    slide_3_overview(prs)
    slide_4_problem_why_now(prs)
    slide_5_solution(prs)
    slide_6_product(prs)
    slide_7_closed_loop(prs)
    slide_8_market(prs)
    slide_9_competitive(prs)
    slide_10_business_model(prs)
    slide_11_unit_economics(prs)
    slide_12_roadmap(prs)
    slide_13_team(prs)
    slide_14_ask(prs)
    slide_15_thank_you(prs)
    slide_a1_qr(prs)
    slide_a2_architecture(prs)
    slide_a3_compliance(prs)

    OUT.parent.mkdir(exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT}  ·  {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
